# -*- coding: utf-8 -*-
"""Cordon — SIH26171 deck generator."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

C = lambda h: RGBColor.from_string(h)
BG, SURF, SURF2, LINE = C("0E1214"), C("161C1F"), C("1D2528"), C("2D383C")
INK, INK2, INK3 = C("E6ECED"), C("B3C0C3"), C("7E8E92")
CLIENT, SERVER, BOUND = C("4CC5D0"), C("9B93F5"), C("E39A57")
OK, DANGER = C("5FC08F"), C("F08A7E")
CLIENT_BG, SERVER_BG, BOUND_BG = C("11292B"), C("1E1B3A"), C("2E2011")

F, FM = "Segoe UI", "Consolas"
W, H = 13.333, 7.5
M = 0.72
CW = W - 2 * M

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
_n = [0]


def rect(s, l, t, w, h, fill=None, line=None, lw=1.0, shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=None):
    sh = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if adj is not None:
        try:
            sh.adjustments[0] = adj
        except Exception:
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def tb(s, l, t, w, h, blocks, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, b in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = b.get("align", align)
        p.space_after = Pt(b.get("after", 0))
        p.space_before = Pt(b.get("before", 0))
        if b.get("line"):
            p.line_spacing = b["line"]
        r = p.add_run()
        r.text = b["t"]
        f = r.font
        f.name = b.get("font", F)
        f.size = Pt(b.get("sz", 14))
        f.bold = b.get("b", False)
        f.color.rgb = b.get("c", INK2)
        if b.get("sp"):
            from pptx.oxml.ns import qn
            r.font._rPr.set("spc", str(int(b["sp"] * 100)))
    return box


def slide(kicker=None, title=None, metric=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    rect(s, 0, 0, W, H, fill=BG, shape=MSO_SHAPE.RECTANGLE)
    _n[0] += 1
    y = 0.52
    if kicker:
        tb(s, M, y, CW - 3.2, 0.26, [{"t": kicker.upper(), "sz": 10.5, "b": True, "c": BOUND, "sp": 1.8}])
        y += 0.34
    if title:
        tb(s, M, y, CW - 3.2, 0.95, [{"t": title, "sz": 30, "b": True, "c": INK}])
    if metric:
        bw = 2.55
        rect(s, W - M - bw, 0.5, bw, 0.42, fill=BOUND_BG, line=BOUND, lw=1.0, adj=0.28)
        tb(s, W - M - bw, 0.5, bw, 0.42, [{"t": metric, "sz": 11, "b": True, "c": BOUND, "align": PP_ALIGN.CENTER}],
           anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    rect(s, M, H - 0.62, CW, 0.012, fill=LINE, shape=MSO_SHAPE.RECTANGLE)
    tb(s, M, H - 0.5, CW / 2, 0.3, [{"t": "CORDON  ·  SIH26171", "sz": 9, "b": True, "c": INK3, "sp": 1.5}])
    tb(s, W / 2, H - 0.5, CW / 2, 0.3, [{"t": "%02d" % _n[0], "sz": 9, "b": True, "c": INK3, "align": PP_ALIGN.RIGHT}],
       align=PP_ALIGN.RIGHT)
    return s


def card(s, l, t, w, h, head, body, accent=CLIENT, num=None, hs=14.5, bs=11.5):
    rect(s, l, t, w, h, fill=SURF, line=LINE, lw=0.75, adj=0.06)
    rect(s, l, t + 0.16, 0.045, h - 0.32, fill=accent, shape=MSO_SHAPE.RECTANGLE)
    px = l + 0.28
    pw = w - 0.5
    y = t + 0.2
    if num:
        tb(s, px, y, pw, 0.22, [{"t": num, "sz": 9.5, "b": True, "c": accent, "sp": 1.6, "font": FM}])
        y += 0.28
    tb(s, px, y, pw, 0.3, [{"t": head, "sz": hs, "b": True, "c": INK}])
    tb(s, px, y + 0.34, pw, h - (y - t) - 0.5, [{"t": body, "sz": bs, "c": INK2, "line": 1.28}])


def bullets(s, l, t, w, items, sz=13.5, gap=0.52, dot=CLIENT):
    y = t
    for it in items:
        rect(s, l, y + 0.115, 0.1, 0.1, fill=dot, shape=MSO_SHAPE.OVAL)
        if isinstance(it, tuple):
            tb(s, l + 0.28, y, w - 0.28, 0.34, [{"t": it[0], "sz": sz, "b": True, "c": INK}])
            tb(s, l + 0.28, y + 0.3, w - 0.28, 0.5, [{"t": it[1], "sz": sz - 2, "c": INK2, "line": 1.25}])
            y += gap + 0.28
        else:
            tb(s, l + 0.28, y, w - 0.28, 0.4, [{"t": it, "sz": sz, "c": INK2, "line": 1.25}])
            y += gap
    return y


def table(s, l, t, w, cols, rows, widths, hsz=10, rsz=11.5, rh=0.42):
    x = l
    for i, cname in enumerate(cols):
        tb(s, x, t, widths[i], 0.25, [{"t": cname.upper(), "sz": hsz, "b": True, "c": INK3, "sp": 1.4}])
        x += widths[i]
    rect(s, l, t + 0.32, w, 0.012, fill=LINE, shape=MSO_SHAPE.RECTANGLE)
    y = t + 0.46
    for row in rows:
        x = l
        for i, cell in enumerate(row):
            txt, col, bold = (cell if isinstance(cell, tuple) else (cell, INK2, False))
            tb(s, x, y, widths[i] - 0.15, rh, [{"t": txt, "sz": rsz, "b": bold, "c": col, "line": 1.2}])
            x += widths[i]
        y += rh
        rect(s, l, y - 0.08, w, 0.008, fill=C("222B2E"), shape=MSO_SHAPE.RECTANGLE)
    return y


def flowbox(s, l, t, w, h, label, sub=None, fill=SURF, line=CLIENT, tc=INK, sz=11.5):
    rect(s, l, t, w, h, fill=fill, line=line, lw=1.0, adj=0.12)
    if sub:
        tb(s, l + 0.08, t, w - 0.16, h, [
            {"t": label, "sz": sz, "b": True, "c": tc, "align": PP_ALIGN.CENTER, "after": 2},
            {"t": sub, "sz": sz - 2.5, "c": INK3, "align": PP_ALIGN.CENTER, "line": 1.1}],
           anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    else:
        tb(s, l + 0.08, t, w - 0.16, h, [{"t": label, "sz": sz, "b": True, "c": tc, "align": PP_ALIGN.CENTER}],
           anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def arrow(s, l, t, w=0.3, h=0.3, c=INK3):
    tb(s, l, t, w, h, [{"t": "→", "sz": 15, "b": True, "c": c, "align": PP_ALIGN.CENTER}],
       anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════ 01 TITLE
s = prs.slides.add_slide(prs.slide_layouts[6])
rect(s, 0, 0, W, H, fill=BG, shape=MSO_SHAPE.RECTANGLE)
rect(s, 0, 0, 0.09, H, fill=CLIENT, shape=MSO_SHAPE.RECTANGLE)
tb(s, M + 0.2, 1.55, CW, 0.3, [{"t": "SMART INDIA HACKATHON  ·  PROBLEM STATEMENT SIH26171", "sz": 11.5, "b": True,
                                "c": BOUND, "sp": 2.0}])
tb(s, M + 0.2, 2.05, CW, 1.5, [{"t": "Cordon", "sz": 66, "b": True, "c": INK}])
tb(s, M + 0.2, 3.15, 9.2, 1.5, [{"t": "On-device visual perception for light-weight browser agents", "sz": 23,
                                 "c": CLIENT, "line": 1.25}])
tb(s, M + 0.2, 3.95, 8.6, 1.2, [{"t": "A browser extension where a local vision model reads the screen and decides — "
                                      "escalating to a server-side VLM only when it must, and never sending a pixel, "
                                      "a password, or a name when it does.", "sz": 13.5, "c": INK2, "line": 1.4}])
bx = M + 0.2
for label, col, bgc in [("LOCAL VISION", CLIENT, CLIENT_BG), ("PRIVACY FILTER", BOUND, BOUND_BG),
                        ("SERVER VLM", SERVER, SERVER_BG)]:
    wd = 2.1
    rect(s, bx, 5.15, wd, 0.44, fill=bgc, line=col, lw=1.0, adj=0.3)
    tb(s, bx, 5.15, wd, 0.44, [{"t": label, "sz": 10.5, "b": True, "c": col, "align": PP_ALIGN.CENTER, "sp": 1.2}],
       anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    bx += wd + 0.22
tb(s, M + 0.2, 6.35, CW, 0.5, [{"t": "Team  ·  <your team name>          Institution  ·  <your college>",
                               "sz": 12, "c": INK3}])

# ══════════════════════════════════ 02 PROBLEM
s = slide("The problem", "Agents need to see your screen. Servers should not.")
tb(s, M, 1.72, 11.2, 0.6, [{"t": "An agentic pipeline with access to your visual context and screen state can "
                                 "automate genuinely complex workflows. Almost all of them run server-side — which "
                                 "caps what a user is willing to share with them.", "sz": 14.5, "c": INK2,
                            "line": 1.4}])
y = 2.75
card(s, M, y, 3.72, 1.9, "Server-side agents", "To reason about your screen they must receive your screen. "
     "Passwords, faces, IDs, account numbers — all uploaded. That ceiling on trust is what limits the tasks "
     "people will hand to an agent.", DANGER, "THE STATUS QUO")
card(s, M + 3.95, y, 3.72, 1.9, "Fully local agents", "Your machine has far fewer resources than a server and "
     "cannot host a full reasoning pipeline. A model small enough to run in a browser tab cannot plan a "
     "multi-step task on its own.", DANGER, "THE OBVIOUS FIX FAILS")
card(s, M + 7.9, y, 3.72, 1.9, "The bridge", "Keep perception and privacy enforcement on the device. Send only "
     "non-sensitive structure — screen layout, application fields — to the server for the reasoning it is "
     "actually needed for.", OK, "WHAT THE PS ASKS FOR")
rect(s, M, 5.05, CW, 1.05, fill=SURF2, line=LINE, lw=0.75, adj=0.06)
tb(s, M + 0.35, 5.22, CW - 0.7, 0.75, [
    {"t": "WHY NOW", "sz": 9.5, "b": True, "c": BOUND, "sp": 1.6, "after": 4},
    {"t": "WebGPU and WebAssembly, plus ONNX Runtime Web and Transformers.js, now make it practical to run real "
          "vision models inside a browser tab. The client can finally see for itself.", "sz": 13, "c": INK2,
     "line": 1.3}])

# ══════════════════════════════════ 03 SOLUTION
s = slide("Our solution", "Cordon — a privacy boundary with a memory")
tb(s, M, 1.72, 11.3, 0.5, [{"t": "Two ideas carry the entire system.", "sz": 14.5, "c": INK2}])
card(s, M, 2.35, 5.6, 2.3, "The server is an escalation path, not a step",
     "A local Router decides whether the task can be resolved on device. \"Click the button labelled Submit\" needs "
     "no server at all — it executes in milliseconds with zero network. Only ambiguity, planning and language "
     "judgement escalate.\n\nThis is the PS's own wording: the local model reads the screen and takes the decision; "
     "the server is engaged if it is required.", CLIENT, "IDEA 1  ·  LOCAL-FIRST", 15, 11.5)
card(s, M + 5.9, 2.35, 5.72, 2.3, "Substitution, so the server reasons without seeing",
     "The client holds a vault mapping typed handles to real values. The server receives EMAIL_1, reasons about "
     "EMAIL_1, and returns actions written in EMAIL_1. The client resolves it locally at the last moment.\n\n"
     "The PS requires the server be aware of the redaction scheme — this handle grammar is that scheme, versioned "
     "and supplied in the system prompt.", SERVER, "IDEA 2  ·  SAFE REFERENCES", 15, 11.5)
rect(s, M, 4.95, CW, 1.2, fill=SURF, line=LINE, lw=0.75, adj=0.05)
cols = [("srikar@gmail.com", "EMAIL_1"), ("\"Hunter2!\"", "never read at all"),
        ("1920×1080 screenshot", "768px, faces masked, opt-in"), ("<input id=\"user_email\">",
                                                                  "{id: el_12, holds: EMAIL_1}")]
tb(s, M + 0.3, 5.08, 4.0, 0.25, [{"t": "STAYS ON DEVICE", "sz": 9.5, "b": True, "c": CLIENT, "sp": 1.5}])
tb(s, M + 6.1, 5.08, 4.0, 0.25, [{"t": "CROSSES THE BOUNDARY", "sz": 9.5, "b": True, "c": BOUND, "sp": 1.5}])
yy = 5.38
for a, b in cols:
    tb(s, M + 0.3, yy, 5.2, 0.22, [{"t": a, "sz": 10.5, "c": INK2, "font": FM}])
    tb(s, M + 5.6, yy, 0.3, 0.22, [{"t": "→", "sz": 10.5, "c": INK3}])
    tb(s, M + 6.1, yy, 5.2, 0.22, [{"t": b, "sz": 10.5, "c": INK, "font": FM}])
    yy += 0.2

# ══════════════════════════════════ 04 HOW IT WORKS
s = slide("How it works", "Eleven steps, end to end")
steps = [("1", "User gives a task", "\"Log in and download my invoice\""),
         ("2", "Read the screen", "DOM + accessibility tree + screenshot"),
         ("3", "Local ViT understands it", "what each element is and means"),
         ("4", "Can I do this myself?", "if yes → jump to step 10"),
         ("5", "Find everything private", "passwords, emails, names, faces, IDs"),
         ("6", "Hide it", "drop · substitute · mask — reals go to the vault")]
steps2 = [("7", "Verify nothing leaked", "independent gate, can veto the send"),
          ("8", "Send the safe version", "structure + handles only"),
          ("9", "Server reasons and replies", "\"type EMAIL_1 into el_12\""),
          ("10", "Resolve + safety-check", "vault lookup, sink check, grounding"),
          ("11", "Execute in the browser", "then verify it actually worked"),
          ("↺", "Loop until done", "re-analyse only what changed")]
for idx, group in enumerate([steps, steps2]):
    yy = 1.78
    lx = M + idx * 6.0
    for num, head, sub in group:
        acc = BOUND if num in ("4", "7") else (SERVER if num in ("8", "9") else CLIENT)
        rect(s, lx, yy, 0.46, 0.46, fill=SURF2, line=acc, lw=1.0, adj=0.22)
        tb(s, lx, yy, 0.46, 0.46, [{"t": num, "sz": 12, "b": True, "c": acc, "align": PP_ALIGN.CENTER}],
           anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
        tb(s, lx + 0.62, yy + 0.02, 4.9, 0.26, [{"t": head, "sz": 13.5, "b": True, "c": INK}])
        tb(s, lx + 0.62, yy + 0.26, 4.9, 0.26, [{"t": sub, "sz": 11, "c": INK3}])
        yy += 0.72
rect(s, M, 6.28, CW, 0.5, fill=BOUND_BG, line=BOUND, lw=0.75, adj=0.2)
tb(s, M + 0.3, 6.28, CW - 0.6, 0.5, [{"t": "Step 3 is worth 25% of the marks.   Step 4 is why we are fast — if it "
                                           "answers yes, steps 5 to 9 never happen.   Step 7 is the promise.",
                                      "sz": 12, "b": True, "c": BOUND}], anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════ 05 ARCHITECTURE
s = slide("Architecture", "Client · boundary · server")
rect(s, M, 1.72, 5.15, 4.35, fill=CLIENT_BG, line=CLIENT, lw=1.0, adj=0.04)
tb(s, M + 0.25, 1.88, 4.6, 0.3, [{"t": "ON DEVICE — TRUSTED", "sz": 10, "b": True, "c": CLIENT, "sp": 1.5}])
yy = 2.3
for t_, d_ in [("1 · Capture", "DOM + accessibility tree + screenshot"),
               ("2 · Perceive", "coverage map → detector → ViT-Tiny → ScreenGraph"),
               ("3 · Route", "resolvable locally? most steps stop here"),
               ("4 · Detect PII", "DOM rules · regex+checksum · NER · vision+OCR"),
               ("5 · Redact", "drop · substitute · tight bbox mask · vault"),
               ("6 · Verify", "independent gate — V1 to V6, has a veto")]:
    tb(s, M + 0.25, yy, 4.7, 0.24, [{"t": t_, "sz": 12.5, "b": True, "c": INK}])
    tb(s, M + 0.25, yy + 0.24, 4.7, 0.24, [{"t": d_, "sz": 10.5, "c": INK3}])
    yy += 0.6
rect(s, M + 5.42, 1.72, 1.05, 4.35, fill=BOUND_BG, line=BOUND, lw=1.25, adj=0.04)
tb(s, M + 5.42, 1.72, 1.05, 4.35, [
    {"t": "PRIVACY", "sz": 10, "b": True, "c": BOUND, "align": PP_ALIGN.CENTER, "after": 2},
    {"t": "BOUNDARY", "sz": 10, "b": True, "c": BOUND, "align": PP_ALIGN.CENTER, "after": 14},
    {"t": "→", "sz": 20, "b": True, "c": BOUND, "align": PP_ALIGN.CENTER, "after": 2},
    {"t": "structure\n+ handles", "sz": 9.5, "c": BOUND, "align": PP_ALIGN.CENTER, "after": 14, "line": 1.1},
    {"t": "←", "sz": 20, "b": True, "c": BOUND, "align": PP_ALIGN.CENTER, "after": 2},
    {"t": "one action\nin handles", "sz": 9.5, "c": BOUND, "align": PP_ALIGN.CENTER, "line": 1.1}],
   anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
rect(s, M + 6.72, 1.72, 4.9, 4.35, fill=SERVER_BG, line=SERVER, lw=1.0, adj=0.04)
tb(s, M + 6.97, 1.88, 4.4, 0.3, [{"t": "SERVER — NEVER TRUSTED WITH RAW DATA", "sz": 10, "b": True, "c": SERVER,
                                  "sp": 1.2}])
yy = 2.3
for t_, d_ in [("7 · Reason", "open-weights VLM, redaction grammar in prompt"),
               ("8 · Constrain", "guided JSON decoding, schema-enforced"),
               ("9 · Guard output", "reject any literal that looks like PII")]:
    tb(s, M + 6.97, yy, 4.4, 0.24, [{"t": t_, "sz": 12.5, "b": True, "c": INK}])
    tb(s, M + 6.97, yy + 0.24, 4.4, 0.24, [{"t": d_, "sz": 10.5, "c": INK3}])
    yy += 0.6
rect(s, M + 6.97, 4.22, 4.4, 0.012, fill=SERVER, shape=MSO_SHAPE.RECTANGLE)
tb(s, M + 6.97, 4.38, 4.4, 1.5, [
    {"t": "RETURNS ONE OF FOUR", "sz": 9.5, "b": True, "c": SERVER, "sp": 1.4, "after": 6},
    {"t": "action   →  click / fill / scroll / navigate", "sz": 11, "c": INK2, "font": FM, "after": 3},
    {"t": "data     →  processed answer for the client", "sz": 11, "c": INK2, "font": FM, "after": 3},
    {"t": "plan     →  up to 3 steps to run locally", "sz": 11, "c": INK2, "font": FM, "after": 3},
    {"t": "ask_user →  ambiguity the user must settle", "sz": 11, "c": INK2, "font": FM}])
rect(s, M, 6.22, CW, 0.55, fill=SURF, line=LINE, lw=0.75, adj=0.18)
tb(s, M + 0.3, 6.22, CW - 0.6, 0.55, [{"t": "Back on device:  resolve handle  →  sink + policy check  →  ground "
                                            "against live DOM  →  execute  →  verify post-condition  →  loop",
                                       "sz": 12, "b": True, "c": CLIENT}], anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════ 06 LOCAL VISION
s = slide("Client component 1", "Local vision processing", "METRIC 1  ·  25%")
tb(s, M, 1.78, 7.6, 0.5, [{"t": "The PS's largest metric is accuracy of visual context. Our output is not a caption "
                                "of the screen — it is a structured, grounded ScreenGraph.", "sz": 14, "c": INK2,
                           "line": 1.35}])
rect(s, M, 2.48, 5.4, 1.5, fill=SURF2, line=LINE, lw=0.75, adj=0.06)
tb(s, M + 0.28, 2.6, 5.0, 1.3, [
    {"t": "SCREENGRAPH", "sz": 9.5, "b": True, "c": CLIENT, "sp": 1.5, "after": 5},
    {"t": "elements[]     id · role · name · bbox · state · conf", "sz": 10.5, "c": INK2, "font": FM, "after": 3},
    {"t": "groups[]       form | nav | modal | table | card", "sz": 10.5, "c": INK2, "font": FM, "after": 3},
    {"t": "reading_order  focus  viewport", "sz": 10.5, "c": INK2, "font": FM}])
card(s, M + 5.7, 2.48, 2.85, 1.5, "DOM + AX tree", "Role, name, type, state. Near-certain, ~0 ms. Blind to canvas, "
     "images, custom widgets.", CLIENT, "CHANNEL A", 12.5, 10.5)
card(s, M + 8.75, 2.48, 2.85, 1.5, "Detector + ViT", "Sees what is actually rendered. Authoritative inside canvas "
     "and images.", CLIENT, "CHANNEL B", 12.5, 10.5)
tb(s, M, 4.18, 6.0, 0.3, [{"t": "FUSION — MATCH BY IoU, THEN ARBITRATE", "sz": 9.5, "b": True, "c": BOUND,
                           "sp": 1.5}])
yy = 4.52
for a, b, col in [("Both agree", "confidence 1.0 — the common case", INK2),
                  ("DOM only, no pixels", "visible: false — occluded or offscreen. Stops the agent clicking a node "
                                          "that is not on screen.", BOUND),
                  ("Vision only, no DOM", "virtual element — canvas app or custom widget. Vision is authoritative.",
                   SERVER),
                  ("Role conflict", "DOM wins for input/button/link; vision wins inside canvas or image.", INK2)]:
    tb(s, M, yy, 2.9, 0.24, [{"t": a, "sz": 12, "b": True, "c": col}])
    tb(s, M + 3.0, yy, 8.6, 0.24, [{"t": b, "sz": 11.5, "c": INK2}])
    yy += 0.42
rect(s, M, 6.26, CW, 0.5, fill=SURF, line=OK, lw=0.75, adj=0.2)
tb(s, M + 0.3, 6.26, CW - 0.6, 0.5, [{"t": "MEASURED, NOT CLAIMED:   element P/R/F1 at IoU ≥ 0.5   ·   role accuracy   "
                                           "·   mean bbox IoU   ·   task-relevant hit rate", "sz": 11.5, "b": True,
                                      "c": OK}], anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════ 07 COVERAGE-GUIDED
s = slide("Client component 1", "Coverage-guided vision — why we stay light", "METRIC 4  ·  20%")
tb(s, M, 1.78, 11.2, 0.5, [{"t": "Running a vision model over every frame loses the resource metric outright. So we "
                                 "run it only where the DOM cannot already answer the question.", "sz": 14,
                            "c": INK2, "line": 1.35}])
flowbox(s, M, 2.55, 2.3, 0.75, "Viewport", "current frame", SURF, LINE, INK)
arrow(s, M + 2.4, 2.72)
flowbox(s, M + 2.8, 2.55, 2.6, 0.75, "Coverage map", "rasterise known bboxes", SURF, CLIENT, INK)
arrow(s, M + 5.5, 2.72)
flowbox(s, M + 5.9, 2.2, 2.75, 0.75, "Explained pixels", "DOM is authoritative · 0 ms", SURF, CLIENT, INK)
flowbox(s, M + 5.9, 3.1, 2.75, 0.75, "Unexplained pixels", "img · canvas · svg · iframe", SURF, BOUND, INK)
arrow(s, M + 8.75, 3.3)
flowbox(s, M + 9.15, 3.1, 2.45, 0.75, "5–20 crops", "not a full frame", SURF, BOUND, INK)
y = bullets(s, M, 4.25, 5.5, [
    ("Two-stage cascade", "A ~1 MB int8 detector proposes; a ViT-Tiny classifies 224px crops, batched on WebGPU."),
    ("Delta re-perception", "MutationObserver marks dirty subtrees. A step after a click re-analyses what changed."),
], 13, 0.42)
bullets(s, M + 6.0, 4.25, 5.6, [
    ("Caching", "Subtree hash → classification. Crop hash → vision result. Repeat frames cost almost nothing."),
    ("WebGPU with a real fallback", "ONNX Runtime Web on WebGPU; WASM + SIMD + threads where WebGPU is absent."),
], 13, 0.42)
rect(s, M, 6.26, CW, 0.5, fill=BOUND_BG, line=BOUND, lw=0.75, adj=0.2)
tb(s, M + 0.3, 6.26, CW - 0.6, 0.5, [{"t": "This is also the honest answer to \"if you have a ViT, why do you need "
                                           "the DOM?\" — they answer different questions, at different costs.",
                                      "sz": 12, "b": True, "c": BOUND}], anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════ 08 PII DETECTION
s = slide("Client component 2", "Privacy filter — detection", "METRIC 2  ·  20%")
tb(s, M, 1.78, 11.2, 0.5, [{"t": "The metric grades recall AND precision. Redacting everything ambiguous would score "
                                 "well on recall and destroy the other half — so detection is calibrated, not "
                                 "paranoid.", "sz": 14, "c": INK2, "line": 1.35}])
table(s, M, 2.5, 11.5, ["Detector", "Finds", "Cost", "Confidence"],
      [[("DOM & AX rules", INK, True), "type=password · autocomplete tokens · label and aria keyword match",
        ("~0 ms", INK2, False), ("0.95 – 1.0", CLIENT, True)],
       [("Regex + checksum", INK, True), "email · phone · Aadhaar (Verhoeff) · PAN · card (Luhn) · IFSC · UPI · JWT",
        ("~1 ms", INK2, False), ("0.55 – 0.98", CLIENT, True)],
       [("Local NER", INK, True), "PERSON · LOCATION · ORG in free text the rules missed",
        ("~15 ms/chunk", INK2, False), ("0.5 – 0.9", CLIENT, True)],
       [("Vision + OCR", INK, True), "faces · ID cards · signatures · QR · text baked into images",
        ("~30 ms/step", INK2, False), ("0.5 – 0.95", CLIENT, True)]],
      [2.5, 5.9, 1.6, 1.5], rh=0.46)
rect(s, M, 4.68, 5.5, 1.45, fill=SURF2, line=LINE, lw=0.75, adj=0.06)
tb(s, M + 0.28, 4.82, 5.0, 1.2, [
    {"t": "FUSION — NOISY-OR, TWO THRESHOLDS", "sz": 9.5, "b": True, "c": BOUND, "sp": 1.3, "after": 6},
    {"t": "p = 1 − Π(1 − pᵢ)", "sz": 12, "b": True, "c": INK, "font": FM, "after": 5},
    {"t": "p ≥ 0.80          → redact", "sz": 10.5, "c": INK2, "font": FM, "after": 2},
    {"t": "0.35 ≤ p < 0.80   → context tie-break", "sz": 10.5, "c": BOUND, "font": FM, "after": 2},
    {"t": "p < 0.35          → keep verbatim", "sz": 10.5, "c": INK2, "font": FM}])
tb(s, M + 5.9, 4.68, 5.7, 0.3, [{"t": "THE TIE-BREAK IS WHERE PRECISION IS WON", "sz": 9.5, "b": True, "c": BOUND,
                                 "sp": 1.3}])
yy = 5.0
for a, b in [("Label proximity", "\"Order number: 1234567890\" vs \"Phone: 9876543210\" — identical to a regex, "
                                 "not to a label"),
             ("Checksum outcome", "12 digits failing Verhoeff is not an Aadhaar. Hard negative evidence."),
             ("Container & repetition", "Inside a form vs. body prose. A value on 50 places is boilerplate.")]:
    tb(s, M + 5.9, yy, 5.7, 0.22, [{"t": a, "sz": 11.5, "b": True, "c": INK}])
    tb(s, M + 5.9, yy + 0.22, 5.7, 0.24, [{"t": b, "sz": 10.5, "c": INK3, "line": 1.15}])
    yy += 0.44

# ══════════════════════════════════ 09 REDACTION
s = slide("Client component 2", "Privacy filter — redaction", "METRIC 3  ·  20%")
tb(s, M, 1.78, 11.2, 0.5, [{"t": "The metric is precision of redaction. Masking pixels that were never sensitive "
                                 "costs marks exactly like missing sensitive ones — so every mask is tight by "
                                 "construction.", "sz": 14, "c": INK2, "line": 1.35}])
table(s, M, 2.5, 11.5, ["Sensitive class", "Method", "What the server sees", "Why"],
      [[("Password · OTP · API key", INK, True), ("Remove completely", DANGER, True), ("sensitive: true", INK2, False),
        "Value never read into the payload — absent, not masked"],
       [("Email · Phone · Name", INK, True), ("Semantic obfuscation", CLIENT, True), ("EMAIL_1  PHONE_1  PERSON_2",
                                                                                     INK2, False),
        "Model keeps the meaning of the field without the value"],
       [("Card number", INK, True), ("Semantic obfuscation", CLIENT, True), ("CARD_1 — no last-4", INK2, False),
        "Last-4 is still identifying, so it is not preserved"],
       [("Face · ID card · signature", INK, True), ("Bounding-box masking", SERVER, True), ("solid fill + class + box",
                                                                                           INK2, False),
        "Solid, not blur — blur and pixelation are recoverable"],
       [("Buttons · links · layout", INK, True), ("No redaction", INK3, True), ("verbatim", INK2, False),
        "The agent cannot operate without them"]],
      [2.85, 2.35, 3.1, 3.2], rh=0.5)
rect(s, M, 5.42, 5.6, 1.35, fill=SURF2, line=OK, lw=0.75, adj=0.06)
tb(s, M + 0.28, 5.56, 5.1, 1.1, [
    {"t": "TIGHT BY CONSTRUCTION", "sz": 9.5, "b": True, "c": OK, "sp": 1.4, "after": 5},
    {"t": "Text  →  Range.getClientRects() over the matched substring only. Mask the eleven characters of the "
          "email, not the paragraph.\nRegions  →  the detector's own box + 4px. Never a whole element, never the "
          "frame.", "sz": 11, "c": INK2, "line": 1.28}])
rect(s, M + 6.0, 5.42, 5.6, 1.35, fill=SURF2, line=LINE, lw=0.75, adj=0.06)
tb(s, M + 6.28, 5.56, 5.1, 1.1, [
    {"t": "MEASURED ON A LABELLED SET", "sz": 9.5, "b": True, "c": OK, "sp": 1.4, "after": 5},
    {"t": "Mask IoU against ground truth   ·   over-redaction rate, the share of safe pixels and characters masked   "
          "·   leak rate, the share of sensitive content left exposed.", "sz": 11, "c": INK2, "line": 1.28}])

# ══════════════════════════════════ 10 VERIFIER
s = slide("The gate", "Only anonymized, unidentifiable data is transmitted")
tb(s, M, 1.78, 11.2, 0.5, [{"t": "The redactor could have a bug. So a second component, sharing no code with it, "
                                 "inspects the serialised payload immediately before fetch — and can veto the send.",
                            "sz": 14, "c": INK2, "line": 1.35}])
checks = [("V1", "Re-scan the bytes", "Full regex + checksum battery over the outgoing JSON string itself, not the "
                                      "object the redactor built"),
          ("V2", "Vault cross-check", "No vault plaintext appears anywhere in the payload — including URLs, alt "
                                      "text and class names"),
          ("V3", "Re-detect on the masked image", "Run the face detector on the redacted bitmap. Still finds a "
                                                  "face? A mask failed to paint. Fail."),
          ("V4", "Entropy sweep", "Any high-entropy string over 20 characters that is not a known handle"),
          ("V5", "Deny-by-default serialiser", "Only whitelisted keys reach the wire. A field added later cannot "
                                               "leak by being forgotten."),
          ("V6", "Escalate, then refuse", "Re-redact harder and re-verify, twice. Then abort the step and tell the "
                                          "user why. Never send.")]
yy = 2.45
for i, (code, head, body) in enumerate(checks):
    col = DANGER if code == "V6" else OK
    lx = M + (i % 2) * 5.95
    if i % 2 == 0 and i > 0:
        yy += 0.0
    ty = 2.45 + (i // 2) * 1.12
    rect(s, lx, ty, 0.52, 0.42, fill=SURF2, line=col, lw=1.0, adj=0.22)
    tb(s, lx, ty, 0.52, 0.42, [{"t": code, "sz": 11.5, "b": True, "c": col, "align": PP_ALIGN.CENTER, "font": FM}],
       anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    tb(s, lx + 0.68, ty + 0.01, 5.0, 0.26, [{"t": head, "sz": 13, "b": True, "c": INK}])
    tb(s, lx + 0.68, ty + 0.27, 5.0, 0.6, [{"t": body, "sz": 10.5, "c": INK3, "line": 1.22}])
rect(s, M, 6.06, CW, 0.62, fill=SURF, line=BOUND, lw=0.75, adj=0.16)
tb(s, M + 0.3, 6.06, CW - 0.6, 0.62, [
    {"t": "PRIVACY RECEIPT — emitted per step into the side panel", "sz": 11, "b": True, "c": BOUND, "after": 3},
    {"t": "counts by class · which detector caught what · mask count · payload bytes · verifier version · payload "
          "hash.  An invisible guarantee becomes something a judge watches happen, live.", "sz": 11, "c": INK2}])

# ══════════════════════════════════ 11 SERVER
s = slide("Server-side integration", "Reasoning over sanitized context")
tb(s, M, 1.78, 11.2, 0.5, [{"t": "The server is aware of the redaction scheme — the handle grammar is versioned and "
                                 "supplied in its system prompt, so it can interpret sanitized data and act on it.",
                            "sz": 14, "c": INK2, "line": 1.35}])
rect(s, M, 2.5, 5.6, 2.55, fill=SURF, line=LINE, lw=0.75, adj=0.05)
tb(s, M + 0.28, 2.64, 5.1, 2.3, [
    {"t": "CLIENT → SERVER", "sz": 9.5, "b": True, "c": CLIENT, "sp": 1.4, "after": 6},
    {"t": '{\n  "schema": "cordon/redaction@1",\n  "task": "Download my latest invoice",\n  "elements": [\n'
          '    { "id":"el_12", "role":"textbox",\n      "name":"Email", "holds":"EMAIL_1" },\n'
          '    { "id":"el_13", "type":"password",\n      "sensitive": true },\n'
          '    { "id":"el_14", "role":"button",\n      "name":"Sign in" } ],\n'
          '  "image": null\n}', "sz": 10, "c": INK2, "font": FM, "line": 1.2}])
rect(s, M + 6.0, 2.5, 5.6, 2.55, fill=SURF, line=LINE, lw=0.75, adj=0.05)
tb(s, M + 6.28, 2.64, 5.1, 2.3, [
    {"t": "SERVER → CLIENT", "sz": 9.5, "b": True, "c": SERVER, "sp": 1.4, "after": 6},
    {"t": '{\n  "type": "action",\n  "thought": "Email field is empty;\n              fill before sign-in.",\n'
          '  "action": {\n    "kind": "fill",\n    "target": "el_12",\n    "value": "EMAIL_1"\n  },\n'
          '  "confidence": 0.91\n}', "sz": 10, "c": INK2, "font": FM, "line": 1.2}])
y = 5.25
tb(s, M, y, 5.6, 0.3, [{"t": "MODEL — OFFLINE DEPLOYABLE, PER THE PS", "sz": 9.5, "b": True, "c": SERVER, "sp": 1.3}])
tb(s, M, y + 0.3, 5.6, 0.9, [{"t": "Open-weights VLM — Qwen2.5-VL-7B or Llama-3.2-11B-Vision. Cloud-hosted during "
                                   "SIH, deployable offline with vLLM. Guided JSON decoding enforces the schema.",
                              "sz": 11.5, "c": INK2, "line": 1.28}])
tb(s, M + 6.0, y, 5.6, 0.3, [{"t": "OUTPUT GUARD", "sz": 9.5, "b": True, "c": SERVER, "sp": 1.3}])
tb(s, M + 6.0, y + 0.3, 5.6, 0.9, [{"t": "value must be a handle or a string the server composed itself. Anything "
                                         "resembling real PII in a response is rejected — a prompt-injected page "
                                         "cannot turn the model into an exfiltration channel.", "sz": 11.5,
                                    "c": INK2, "line": 1.28}])

# ══════════════════════════════════ 12 LATENCY/ACCURACY
s = slide("The trade-off the PS asks for", "Latency against accuracy, as a visible control", "METRIC 5  ·  15%")
tb(s, M, 1.78, 11.2, 0.5, [{"t": "\"Participants must balance the trade-offs between inference latency and the "
                                 "accuracy.\"  We make that a mode the user and the judge can switch, not a hidden "
                                 "constant.", "sz": 14, "c": INK2, "line": 1.35}])
table(s, M, 2.55, 11.5, ["Mode", "Vision", "Behaviour", "Typical step"],
      [[("Fast", CLIENT, True), "cached only", "DOM/AX graph, regex + DOM PII, local routing only",
        ("~15 ms", OK, True)],
       [("Balanced  (default)", CLIENT, True), "coverage-guided, int8",
        "+ ViT on unexplained regions, NER on flagged text", ("~90 ms", OK, True)],
       [("Thorough", CLIENT, True), "full frame", "+ OCR everywhere, full-frame detector sweep",
        ("~350 ms", BOUND, True)]],
      [2.6, 2.3, 5.3, 1.6], rh=0.52)
y = bullets(s, M, 4.42, 5.5, [
    ("Local-first is the real latency win", "Most steps resolve on device with zero network. The round trip is the "
                                            "dominant cost, so removing it beats optimising it."),
    ("Warm sessions, cached weights", "Models fetched once into OPFS. ORT sessions kept alive. No cold start per "
                                      "step."),
], 13, 0.44)
bullets(s, M + 6.0, 4.42, 5.6, [
    ("One action per round trip", "Safer to ground, and lets the server return a short local plan when it is "
                                  "confident."),
    ("Demoed at all three modes", "The same task run in Fast, Balanced and Thorough, with accuracy numbers beside "
                                  "the latency numbers, is the answer to that requirement."),
], 13, 0.44)

# ══════════════════════════════════ 13 TECH STACK
s = slide("Technology", "Stack — and what is built today")
# (label, status) — 1 built, 0 planned, 2 partial
groups = [("CLIENT — EXTENSION", CLIENT, [("WebExtension Manifest V3", 1), ("Chrome build", 1),
                                          ("Firefox build", 2), ("TypeScript + esbuild", 1),
                                          ("Content script + service worker", 1)]),
          ("ON-DEVICE INFERENCE", CLIENT, [("ONNX Runtime Web", 0), ("WebGPU, WASM + SIMD fallback", 0),
                                           ("ViT-Tiny crop classifier", 0), ("int8 quantised UI detector", 0),
                                           ("OffscreenCanvas + Web Workers", 0)]),
          ("BROWSER APIS", BOUND, [("DOM + accessible names", 1), ("captureVisibleTab", 0),
                                   ("Mutation / Intersection observers", 0), ("Range.getClientRects", 1),
                                   ("Canvas bounding-box masking", 0)]),
          ("SERVER", SERVER, [("Rule-based planner, Node", 1), ("Open-weights VLM via vLLM", 0),
                              ("Qwen2.5-VL / Llama-3.2-Vision", 0), ("Guided JSON decoding", 0),
                              ("PII output guard", 1)])]
for i, (name, col, items) in enumerate(groups):
    lx = M + i * 2.98
    rect(s, lx, 1.85, 2.78, 3.5, fill=SURF, line=LINE, lw=0.75, adj=0.05)
    rect(s, lx, 1.85, 2.78, 0.055, fill=col, shape=MSO_SHAPE.RECTANGLE)
    tb(s, lx + 0.24, 2.08, 2.4, 0.4, [{"t": name, "sz": 9.5, "b": True, "c": col, "sp": 1.3, "line": 1.15}])
    yy = 2.62
    for it, built in items:
        mark, mc = ("BUILT", OK) if built == 1 else (("PART", BOUND) if built == 2 else ("PHASE 2-4", INK3))
        tb(s, lx + 0.24, yy, 1.75, 0.45, [{"t": it, "sz": 10.5,
                                           "c": INK2 if built else INK3, "line": 1.18}])
        tb(s, lx + 2.0, yy + 0.02, 0.62, 0.24, [{"t": mark, "sz": 7.5, "b": True, "c": mc,
                                                 "align": PP_ALIGN.RIGHT, "sp": 0.8}], align=PP_ALIGN.RIGHT)
        yy += 0.5
rect(s, M, 5.6, CW, 1.1, fill=SURF2, line=BOUND, lw=0.9, adj=0.06)
tb(s, M + 0.32, 5.76, CW - 0.64, 0.85, [
    {"t": "WHERE WE ARE", "sz": 9.5, "b": True, "c": BOUND, "sp": 1.4, "after": 5},
    {"t": "Phase 1 was deliberately the complete agent loop with no models in it — perception, detection, "
          "redaction, verification, routing and grounded execution, end to end, so the vision channel drops into "
          "a working system rather than a scaffold. The on-device inference column is phase 2 and is not yet "
          "implemented.", "sz": 11.5, "c": INK2, "line": 1.28}])

# ══════════════════════════════════ 14 DEMO
s = slide("The prototype", "One end-to-end task, done well")
rect(s, M, 1.8, CW, 0.72, fill=BOUND_BG, line=BOUND, lw=1.0, adj=0.16)
tb(s, M + 0.35, 1.8, CW - 0.7, 0.72, [{"t": "\u201cFill this job application from my profile.\u201d",
                                       "sz": 19, "b": True, "c": BOUND}], anchor=MSO_ANCHOR.MIDDLE)
tb(s, M, 2.72, 11.3, 0.4, [{"t": "One task chosen because it exercises every subsystem the PS names — and every "
                                 "metric it grades.", "sz": 13.5, "c": INK2}])
items = [("Multi-section form", "name, email, phone, experience textarea, file upload", CLIENT),
         ("A photo and an ID-card image on the page", "the DOM sees only <img> — vision is the only channel that "
                                                      "can tell these are sensitive", BOUND),
         ("PII across both channels", "text PII in the DOM, visual PII in the images, OCR'd text inside the ID card",
          BOUND),
         ("Handles resolved locally", "the server plans which profile slot fills which field, never learning a "
                                      "single value", SERVER),
         ("Grounding under re-render", "the form re-renders between steps — the agent re-perceives instead of "
                                       "clicking a stale id", CLIENT),
         ("Human confirmation before submit", "an irreversible action always stops for explicit approval", DANGER)]
yy = 3.28
for head, body, col in items:
    rect(s, M, yy + 0.09, 0.12, 0.12, fill=col, shape=MSO_SHAPE.OVAL)
    tb(s, M + 0.32, yy, 4.6, 0.26, [{"t": head, "sz": 12.5, "b": True, "c": INK}])
    tb(s, M + 5.0, yy, 6.6, 0.32, [{"t": body, "sz": 11.5, "c": INK2, "line": 1.2}])
    yy += 0.51

# ══════════════════════════════════ 14b CURRENT STATE
s = slide("Status", "What runs today, and what does not")
tb(s, M, 1.8, 11.3, 0.4, [{"t": "3,633 lines of TypeScript. Four dependencies, all build tooling. The extension "
                                "and server ship with zero third-party runtime code.", "sz": 13.5, "c": INK2}])
table(s, M, 2.35, 11.5, ["Capability", "State", "Evidence"],
      [[("Screen understanding — DOM + accessible names", INK, True), ("BUILT", OK, True),
        "Live overlay labels every control; occlusion and offscreen distinguished"],
       [("Screen understanding — vision / ViT", INK, True), ("PHASE 2", INK3, True),
        "No pixel is captured yet; the visual half of metric 1 is not addressed"],
       [("PII detection — DOM rules + regex/checksums", INK, True), ("BUILT", OK, True),
        "P 1.000 / R 1.000 on the fixture, with Indian hard negatives"],
       [("PII detection — local NER, OCR", INK, True), ("PHASE 3", INK3, True), "2 of the 4 detectors exist"],
       [("Redaction — text spans, typed handles", INK, True), ("BUILT", OK, True),
        "Character-offset substitution, not element blanking"],
       [("Redaction — image bounding boxes", INK, True), ("PHASE 2", INK3, True),
        "Needs the capture and canvas pipeline"],
       [("Privacy verifier V1 V2 V4 V5 + V6 escalation", INK, True), ("BUILT", OK, True),
        "Independent of the redactor; holds a veto; refuses rather than sends"],
       [("Privacy verifier V3 — re-detect on masked image", INK, True), ("STUB", BOUND, True),
        "Honest no-op until there is an image to re-scan"],
       [("Local-first router", INK, True), ("BUILT", OK, True), "5 of 9 benchmark tasks never touch the network"],
       [("Encrypted profile vault — AES-256-GCM", INK, True), ("BUILT", OK, True),
        "Beyond the original plan; PBKDF2 310k, key in session memory"],
       [("Server — open-weights VLM", INK, True), ("PHASE 4", INK3, True),
        "Rule-based planner speaks the same contract; one function to swap"]],
      [4.6, 1.35, 5.55], rh=0.4, rsz=10.5)

# ══════════════════════════════════ 15 EVALUATION
s = slide("Evaluation", "Three of the five metrics are measured quantities — so we measure them")
tb(s, M, 1.82, 11.2, 0.4, [{"t": "A labelled evaluation harness ships in phase one, not as an afterthought. Judges "
                                 "score what you can put a number on.", "sz": 14, "c": INK2}])
sets = [("eval/screens/", CLIENT, "40–60 labelled pages — real sites saved offline plus our demo pages. Ground "
                                  "truth: every interactive element's bbox, role and accessible name.",
         "→ element P/R/F1 · role accuracy · mean bbox IoU · task-relevant hit rate"),
        ("eval/pii/", CLIENT, "Indian-context positives — Aadhaar, PAN, UPI, IFSC, IN phone formats. Hard negatives "
                              "— order IDs, tracking numbers, PIN codes, prices, dates.",
         "→ precision, recall, F1 per class · thresholds tuned on this set"),
        ("eval/redaction/", CLIENT, "Ground-truth sensitive regions per screenshot, for both text spans and image "
                                    "regions.",
         "→ mask IoU · over-redaction rate · leak rate")]
yy = 2.42
for name, col, body, out in sets:
    rect(s, M, yy, CW, 1.18, fill=SURF, line=LINE, lw=0.75, adj=0.08)
    rect(s, M, yy + 0.14, 0.045, 0.9, fill=col, shape=MSO_SHAPE.RECTANGLE)
    tb(s, M + 0.3, yy + 0.16, 2.7, 0.3, [{"t": name, "sz": 13, "b": True, "c": INK, "font": FM}])
    tb(s, M + 3.15, yy + 0.14, 5.2, 0.9, [{"t": body, "sz": 11, "c": INK2, "line": 1.25}])
    tb(s, M + 8.6, yy + 0.14, 3.0, 0.9, [{"t": out, "sz": 10.5, "b": True, "c": OK, "line": 1.25}])
    yy += 1.3
rect(s, M, 6.32, CW, 0.44, fill=SURF2, line=LINE, lw=0.75, adj=0.22)
tb(s, M + 0.3, 6.32, CW - 0.6, 0.44, [{"t": "eval/run.ts  prints the scorecard.  Metrics 4 and 5 come live from the "
                                            "telemetry HUD: per-stage ms, memory, payload bytes, p50/p95 per task.",
                                       "sz": 11.5, "c": INK2}], anchor=MSO_ANCHOR.MIDDLE)

# ══════════════════════════════════ 16 METRIC MAP
s = slide("Traceability", "Every metric, and where it is earned")
table(s, M, 1.95, 11.5, ["#", "Evaluation metric", "Weight", "Where it is earned", "Evidence in the demo"],
      [[("1", BOUND, True), ("Accuracy of visual context from screen", INK, True), ("25%", BOUND, True),
        "ScreenGraph — DOM + AX fused with detector and ViT-Tiny, IoU-matched and arbitrated",
        "Live overlay draws labelled boxes on any page"],
       [("2", BOUND, True), ("Recall and precision for PII detection", INK, True), ("20%", BOUND, True),
        "Four detectors, noisy-OR, two thresholds, context tie-break with checksums",
        "P/R/F table from the labelled corpus"],
       [("3", BOUND, True), ("Precision of redaction", INK, True), ("20%", BOUND, True),
        "Per-span getClientRects masks, detector box + 4px, never whole elements",
        "Side-by-side raw vs. sanitized payload viewer"],
       [("4", BOUND, True), ("Client-side resource utilization", INK, True), ("20%", BOUND, True),
        "Coverage-guided vision, delta re-perception, caching, int8 models, local-first routing",
        "Telemetry HUD — per-stage ms and MB"],
       [("5", BOUND, True), ("Overall end-to-end latency", INK, True), ("15%", BOUND, True),
        "Most steps never reach the network; warm sessions; one action per round trip",
        "Step timer, p50/p95, three latency modes"]],
      [0.42, 3.15, 0.85, 4.05, 3.03], rh=0.72, rsz=11)

# ══════════════════════════════════ 17 ROADMAP
s = slide("Delivery", "Build order — each phase demonstrable on its own")
phases = [("PHASE 1", "The loop, without AI", "WebExtension MV3 · ScreenGraph from DOM + AX · DOM and regex+checksum "
                                              "detectors · calibrated fusion · handles and vault · verifier · Router "
                                              "with local execution · mock server · executor with grounding · "
                                              "telemetry HUD · the three eval harnesses · demo pages", OK,
           "A complete, honest end-to-end agent"),
          ("PHASE 2", "Pixels, masks, vision channel", "Screenshot → OffscreenCanvas → coverage map · ONNX Runtime "
                                                       "Web with WebGPU and WASM fallback · face and ID detector · "
                                                       "tight canvas bbox masking · verifier V3 · fused ScreenGraph",
           CLIENT, "Metrics 1 and 3 become real"),
          ("PHASE 3", "The ViT and the reader", "UI element detector fine-tuned on a UI corpus · ViT-Tiny crop "
                                                "classifier · OCR on text-bearing crops routed back through the "
                                                "text detectors · local NER · the three latency modes", CLIENT,
           "Full on-device perception"),
          ("PHASE 4", "Real reasoning", "Open-weights VLM behind guided JSON decoding · redaction grammar in the "
                                        "system prompt · server-side PII output guard · Firefox build · final "
                                        "evaluation numbers", SERVER, "End-to-end task, both browsers")]
yy = 1.9
for tag, head, body, col, note in phases:
    rect(s, M, yy, CW, 1.07, fill=SURF, line=LINE, lw=0.75, adj=0.08)
    rect(s, M, yy + 0.13, 0.045, 0.81, fill=col, shape=MSO_SHAPE.RECTANGLE)
    tb(s, M + 0.3, yy + 0.15, 1.5, 0.25, [{"t": tag, "sz": 10, "b": True, "c": col, "sp": 1.4, "font": FM}])
    tb(s, M + 0.3, yy + 0.42, 2.6, 0.5, [{"t": head, "sz": 13, "b": True, "c": INK, "line": 1.15}])
    tb(s, M + 3.15, yy + 0.15, 6.1, 0.85, [{"t": body, "sz": 10.5, "c": INK2, "line": 1.25}])
    tb(s, M + 9.5, yy + 0.15, 2.1, 0.85, [{"t": note, "sz": 10.5, "b": True, "c": col, "line": 1.25}])
    yy += 1.19

# ══════════════════════════════════ 18 FEASIBILITY
s = slide("Feasibility and viability", "The risks we can name, and what we do about them")
risks = [("Vision model too slow in a browser tab", "Coverage-guided scheduling means 5–20 small crops, not a full "
                                                    "frame. int8 quantisation, batched WebGPU inference, cached "
                                                    "results by crop hash.", CLIENT),
         ("WebGPU unavailable on the judging machine", "ONNX Runtime Web falls back to WASM + SIMD + threads. The "
                                                       "Fast mode runs with no vision at all and still completes "
                                                       "the task.", CLIENT),
         ("Over-redaction destroys the precision metrics", "Two thresholds and a context tie-break, tuned on a "
                                                           "labelled corpus with deliberate hard negatives.", BOUND),
         ("Agent clicks the wrong element after a re-render", "Stability signatures are re-checked before every "
                                                              "action; a mismatch triggers re-perception instead of "
                                                              "a click.", BOUND),
         ("Prompt injection from a hostile page", "The server can only return handles; allowedSinks constrains where "
                                                  "each handle may be written; irreversible actions always need a "
                                                  "human.", DANGER),
         ("Firefox MV3 differences", "Browser-specific code isolated behind src/platform/ with webextension-polyfill "
                                     "and two manifests.", SERVER)]
for i, (head, body, col) in enumerate(risks):
    lx = M + (i % 2) * 5.95
    ty = 1.95 + (i // 2) * 1.55
    rect(s, lx, ty, 5.65, 1.36, fill=SURF, line=LINE, lw=0.75, adj=0.07)
    rect(s, lx, ty + 0.15, 0.045, 1.06, fill=col, shape=MSO_SHAPE.RECTANGLE)
    tb(s, lx + 0.3, ty + 0.18, 5.1, 0.5, [{"t": head, "sz": 12.5, "b": True, "c": INK, "line": 1.15}])
    tb(s, lx + 0.3, ty + 0.66, 5.1, 0.62, [{"t": body, "sz": 10.5, "c": INK2, "line": 1.25}])

# ══════════════════════════════════ 19 IMPACT
s = slide("Impact", "What this unlocks")
y = bullets(s, M, 1.95, 5.5, [
    ("A higher ceiling on trust", "The tasks people refuse to hand an agent today are refused because of what would "
                                  "have to be uploaded. Remove the upload and the refusal goes with it."),
    ("Regulation-friendly by construction", "Data minimisation is not a policy here, it is the transport format. "
                                            "The server cannot leak what it never received."),
    ("Cheaper to operate at scale", "Perception, detection and redaction run on a million clients. The fleet only "
                                    "sees the hard fraction, as ~10 KB of sanitized JSON."),
], 13.5, 0.5)
bullets(s, M + 6.0, 1.95, 5.6, [
    ("Portable beyond the browser", "The boundary — perceive locally, substitute, verify, escalate — is the same "
                                    "design for a desktop or mobile agent."),
    ("Auditable privacy", "The per-step receipt makes the guarantee inspectable by the user, not just asserted in a "
                          "policy document."),
    ("Future work", "Federated tuning of the detectors without collecting screens; a shared redaction-scheme spec so "
                    "any server can be made scheme-aware."),
], 13.5, 0.5)
rect(s, M, 5.35, CW, 1.35, fill=SURF2, line=BOUND, lw=1.0, adj=0.06)
tb(s, M + 0.4, 5.52, CW - 0.8, 1.05, [
    {"t": "IN ONE LINE", "sz": 9.5, "b": True, "c": BOUND, "sp": 1.5, "after": 6},
    {"t": "The client sees everything and decides most things itself. The server receives a censored map with "
          "placeholders, and only ever answers \u201cdo this to that box.\u201d", "sz": 16, "b": True, "c": INK,
     "line": 1.3}])

prs.save(r"c:\Users\srika\OneDrive\Desktop\sih-p\docs\Cordon_SIH26171.pptx")
print("saved", _n[0] + 1, "slides")
