# -*- coding: utf-8 -*-
"""
SIH 2026 Idea Submission — six slides.

Visual format modelled on the SIH 2025 winning deck the team supplied:
  · title inside a rounded-rectangle outline, team oval top-left, logo top-right
  · two large black-bordered boxes side by side
  · blue diamond section headers, green underlined sub-headers, black bullets
  · a real flowchart / architecture diagram in the right-hand box
  · numbered references with links on the final slide

Template rules still enforced: 6 slides max, points not paragraphs, the idea
detail pointers reproduced verbatim, export to PDF before uploading.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE

C = lambda h: RGBColor.from_string(h)

BLACK = C("000000")
INK = C("1A1A1A")
HDR = C("1F4E79")       # ❖ section headers
GREEN = C("15803D")     # underlined sub-headers
BLUE = C("0563C1")      # hyperlinks
RED = C("A32B20")
CREAM = C("FFF2CC")     # flowchart fill, as in the winning deck
CREAM2 = C("DEEBF7")
GREY = C("F2F2F2")
WHITE = C("FFFFFF")

SERIF = "Times New Roman"
BODY = "Calibri"
MONO = "Consolas"

W, H = 13.333, 7.5

# ── EDIT THESE ───────────────────────────────────────────────────────────────
TEAM_NAME = "Your Team Name"
TEAM_ID = "<Team ID>"
PS_ID = "SIH26171"
PS_TITLE = ("On-device Visual Perception for Light-weight Browser Agents "
            "(privacy-preserving browser agent with local redaction)")
THEME = "Miscellaneous"
CATEGORY = "Software"
# ─────────────────────────────────────────────────────────────────────────────

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(W), Inches(H)
_n = [0]


def rect(s, l, t, w, h, fill=None, line=BLACK, lw=1.0, shape=MSO_SHAPE.RECTANGLE, adj=None):
    sh = s.shapes.add_shape(shape, Inches(l), Inches(t), Inches(w), Inches(h))
    if adj is not None:
        try:
            sh.adjustments[0] = adj
        except Exception:
            pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    return sh


def tb(s, l, t, w, h, blocks, anchor=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    box = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, b in enumerate(blocks):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = b.get("align", align)
        p.space_after = Pt(b.get("after", 0))
        p.space_before = Pt(b.get("before", 0))
        if b.get("line"):
            p.line_spacing = b["line"]
        for seg in (b["t"] if isinstance(b["t"], list) else [b["t"]]):
            txt, over = seg if isinstance(seg, tuple) else (seg, {})
            r = p.add_run()
            r.text = txt
            f = r.font
            f.name = over.get("font", b.get("font", BODY))
            f.size = Pt(over.get("sz", b.get("sz", 11)))
            f.bold = over.get("b", b.get("b", False))
            f.underline = over.get("u", b.get("u", False))
            f.color.rgb = over.get("c", b.get("c", INK))
    return box


def arrow(s, x1, y1, x2, y2, lw=1.25):
    cn = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.line.color.rgb = BLACK
    cn.line.width = Pt(lw)
    cn.line._get_or_add_ln().append(
        cn.line._get_or_add_ln().makeelement(
            "{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd",
            {"type": "triangle", "w": "med", "len": "med"}))
    return cn


def fbox(s, l, t, w, h, text, fill=CREAM, sz=9.5, bold=True, col=INK):
    """A flowchart node in the winning deck's style: cream fill, black border."""
    rect(s, l, t, w, h, fill=fill, line=BLACK, lw=1.0)
    tb(s, l + 0.04, t, w - 0.08, h, [{"t": text, "sz": sz, "b": bold, "c": col,
                                      "align": PP_ALIGN.CENTER, "line": 0.95}],
       anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)


def slide(title):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _n[0] += 1
    # team oval, top left
    rect(s, 0.18, 0.12, 1.1, 0.62, fill=None, line=C("7030A0"), lw=1.25, shape=MSO_SHAPE.OVAL)
    tb(s, 0.18, 0.12, 1.1, 0.62, [{"t": TEAM_NAME, "sz": 9, "c": INK,
                                   "align": PP_ALIGN.CENTER, "line": 0.9}],
       anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    # title inside a rounded-rectangle outline
    rect(s, 3.95, 0.1, 5.45, 0.66, fill=None, line=BLACK, lw=1.5,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, adj=0.4)
    tb(s, 3.95, 0.1, 5.45, 0.66, [{"t": title, "sz": 23, "b": True, "c": BLACK,
                                   "font": SERIF, "align": PP_ALIGN.CENTER}],
       anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    # logo placeholder
    rect(s, W - 2.15, 0.1, 1.95, 0.66, fill=None, line=C("BFBFBF"), lw=0.75)
    tb(s, W - 2.15, 0.1, 1.95, 0.66, [{"t": "paste SIH 2026 logo", "sz": 9, "c": C("808080"),
                                       "align": PP_ALIGN.CENTER}],
       anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    # page number
    tb(s, W - 0.75, H - 0.42, 0.45, 0.3, [{"t": str(_n[0]), "sz": 12, "b": True, "c": BLACK,
                                           "align": PP_ALIGN.RIGHT}], align=PP_ALIGN.RIGHT)
    return s


def diamond_header(s, l, t, w, text, sz=19):
    tb(s, l, t, w, 0.4, [{"t": [("\u2756 ", {"sz": sz, "b": True, "c": HDR}),
                                (text, {"sz": sz, "b": True, "c": HDR, "u": True})]}])


def subhead(s, l, t, w, text, sz=11):
    tb(s, l, t, w, 0.24, [{"t": text, "sz": sz, "b": True, "c": GREEN, "u": True}])


def bullets(s, l, t, w, items, sz=10.5, gap=0.235, indent=0.16):
    y = t
    for it in items:
        tb(s, l, y, w, 0.3, [{"t": [("\u2022  ", {"c": INK}), (it, {})], "sz": sz, "line": 0.98}])
        y += gap
    return y


# ═══════════════════════════════════════════════════════════ 1 · TITLE
s = prs.slides.add_slide(prs.slide_layouts[6])
_n[0] = 1
tb(s, 0.4, 0.28, 9.6, 0.8, [{"t": "SMART INDIA HACKATHON 2026", "sz": 36, "b": True,
                             "c": HDR, "font": SERIF, "align": PP_ALIGN.CENTER}],
   anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
rect(s, W - 2.75, 0.2, 2.5, 1.0, fill=None, line=C("BFBFBF"), lw=0.75)
tb(s, W - 2.75, 0.2, 2.5, 1.0, [{"t": "paste SIH 2026 logo", "sz": 10, "c": C("808080"),
                                 "align": PP_ALIGN.CENTER}],
   anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
tb(s, 0, 1.3, W, 0.6, [{"t": "TITLE PAGE", "sz": 27, "b": True, "c": BLACK, "font": SERIF,
                        "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)

y = 2.35
for k, v in [("Problem Statement ID \u2013", PS_ID),
             ("Problem Statement Title-", PS_TITLE),
             ("Theme-", THEME),
             ("PS Category- ", CATEGORY),
             ("Team ID- ", TEAM_ID),
             ("Team Name- ", TEAM_NAME)]:
    tb(s, 0.85, y, 11.6, 0.62, [{"t": [("\u2022  ", {"sz": 15, "b": True, "c": BLACK}),
                                       (k, {"sz": 15, "b": True, "c": BLACK}),
                                       (v, {"sz": 15, "c": INK, "u": True})], "line": 1.0}])
    y += 0.62 if len(v) < 60 else 0.86

tb(s, W - 0.75, H - 0.42, 0.45, 0.3, [{"t": "1", "sz": 12, "b": True, "c": BLACK,
                                       "align": PP_ALIGN.RIGHT}], align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════ 2 · IDEA TITLE
s = slide("IDEA TITLE")
LB, RB = 0.3, 6.95          # left box x, right box x
LW, RW = 6.4, 6.05

rect(s, LB, 0.9, LW, 6.25, fill=WHITE, line=BLACK, lw=1.25)
rect(s, RB, 0.9, RW, 6.25, fill=WHITE, line=BLACK, lw=1.25)

diamond_header(s, LB + 0.18, 1.0, LW - 0.3, "Proposed Solution:")

tb(s, LB + 0.18, 1.5, LW - 0.35, 0.5,
   [{"t": [("Problem:  ", {"b": True, "c": RED}),
           ("every browser AI agent today uploads your screen \u2014 passwords, faces and all.",
            {"c": INK})], "sz": 10.5, "line": 1.0}])

items = [
    ("On-device Screen Reading:",
     "A 1.2 MB vision model reads the page inside your browser."),
    ("Local Privacy Engine:",
     "DOM rules, checksums and pixels find every personal value."),
    ("Safe Reference Substitution:",
     "Your email becomes EMAIL_1 \u2014 typed, opaque, and stable."),
    ("Verify-Before-Send Gate:",
     "An independent check re-reads the bytes and can refuse to send."),
    ("Server Reasons Blindly:",
     "It plans using EMAIL_1 and never learns the address."),
    ("Local Resolution & Execution:",
     "Your browser swaps the handle back and types the real value."),
    ("Encrypted Personal Vault:",
     "AES-256-GCM; the key dies when the browser closes."),
    ("Ingestion Verification:",
     "Every field is read back to confirm the value truly landed."),
]
y = 2.06
for head, body in items:
    subhead(s, LB + 0.18, y, LW - 0.35, "\u2022 " + head, sz=10.5)
    tb(s, LB + 0.34, y + 0.21, LW - 0.5, 0.24, [{"t": body, "sz": 10, "c": INK, "line": 0.98}])
    y += 0.47

rect(s, LB + 0.18, 5.92, LW - 0.36, 1.06, fill=C("E9F2EC"), line=GREEN, lw=1.0)
tb(s, LB + 0.32, 6.02, LW - 0.62, 0.9, [
    {"t": "Innovation & Uniqueness \u2014 built and measured today:", "sz": 10, "b": True,
     "c": GREEN, "after": 3},
    {"t": "\u2022 Handles keep meaning: server sees two fields share ONE email, never the email.",
     "sz": 9.5, "c": INK, "line": 0.98, "after": 2},
    {"t": "\u2022 Vision runs only where the page is opaque \u2014 92% of a screen is never analysed.",
     "sz": 9.5, "c": INK, "line": 0.98, "after": 2},
    {"t": "\u2022 A verifier with a veto \u2014 checks the outgoing bytes AND the masked picture.",
     "sz": 9.5, "c": INK, "line": 0.98}])

# ── right box: the flowchart ────────────────────────────────────────────────
cx = RB + RW / 2
fbox(s, cx - 1.5, 1.06, 3.0, 0.38, "USER TASK  \u2014  \u201cfill this form\u201d", CREAM2, 10)
arrow(s, cx, 1.44, cx, 1.62)

rect(s, RB + 0.18, 1.62, RW - 0.36, 2.32, fill=C("F7FBFF"), line=C("9DC3E6"), lw=1.0)
tb(s, RB + 0.28, 1.68, 3.0, 0.22, [{"t": "ON YOUR DEVICE", "sz": 8.5, "b": True, "c": HDR}])

fbox(s, cx - 1.5, 1.94, 3.0, 0.36, "Capture DOM + accessibility tree + screen", CREAM, 9)
arrow(s, cx, 2.30, cx, 2.44)
fbox(s, cx - 1.5, 2.44, 3.0, 0.36, "Vision model reads the pixels", CREAM, 9)
arrow(s, cx, 2.80, cx, 2.94)
fbox(s, cx - 1.5, 2.94, 3.0, 0.36, "Find PII  \u00b7  redact  \u00b7  mask faces", CREAM, 9)
arrow(s, cx, 3.30, cx, 3.44)
fbox(s, cx - 1.5, 3.44, 3.0, 0.40, "VERIFIER  \u2014  may refuse to send", C("FBE5D6"), 9.5)

arrow(s, cx, 3.84, cx, 4.06)
rect(s, RB + 0.18, 4.06, RW - 0.36, 0.34, fill=C("FBEBE9"), line=RED, lw=1.25)
tb(s, RB + 0.18, 4.06, RW - 0.36, 0.34, [{"t": "PRIVACY BOUNDARY  \u2014  handles only, never values",
                                          "sz": 9.5, "b": True, "c": RED, "align": PP_ALIGN.CENTER}],
   anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
arrow(s, cx, 4.40, cx, 4.58)

rect(s, RB + 0.18, 4.58, RW - 0.36, 0.86, fill=C("EDEBF8"), line=C("8E86F0"), lw=1.0)
tb(s, RB + 0.28, 4.64, 3.0, 0.22, [{"t": "SERVER", "sz": 8.5, "b": True, "c": C("4A42B8")}])
fbox(s, cx - 1.5, 4.9, 3.0, 0.46, "Open-weights VLM plans the step\n\u201cfill el_7 with EMAIL_1\u201d",
     C("E4E0F7"), 9)

arrow(s, cx, 5.44, cx, 5.62)
rect(s, RB + 0.18, 5.62, RW - 0.36, 1.36, fill=C("F7FBFF"), line=C("9DC3E6"), lw=1.0)
tb(s, RB + 0.28, 5.68, 3.0, 0.22, [{"t": "BACK ON YOUR DEVICE", "sz": 8.5, "b": True, "c": HDR}])
fbox(s, cx - 1.5, 5.94, 3.0, 0.34, "Resolve EMAIL_1 from the local vault", CREAM, 9)
arrow(s, cx, 6.28, cx, 6.40)
fbox(s, cx - 1.5, 6.40, 3.0, 0.34, "Safety check \u00b7 ground \u00b7 type \u00b7 verify", CREAM, 9)

# ═══════════════════════════════════════════════════════════ 3 · TECHNICAL
s = slide("TECHNICAL APPROACH")
rect(s, LB, 0.9, LW, 6.25, fill=WHITE, line=BLACK, lw=1.25)
rect(s, RB, 0.9, RW, 6.25, fill=WHITE, line=BLACK, lw=1.25)

tb(s, LB + 0.18, 1.0, LW - 0.3, 0.3, [{"t": "\u2022  Technologies to be used (e.g. programming languages, "
                                             "frameworks, hardware)", "sz": 11.5, "b": True, "c": BLACK}])
tech = [("Client Extension:", "TypeScript, WebExtension Manifest V3 (Chrome, Firefox)."),
        ("On-Device Inference:", "ONNX Runtime Web \u2014 WebGPU, falling back to WASM + SIMD."),
        ("Vision Model:", "UltraFace RFB-320 \u2014 1.2 MB CNN, runs inside the browser."),
        ("Screen Understanding:", "DOM + accessibility tree fused with model detections."),
        ("PII Detection:", "DOM rules, regex, Verhoeff & Luhn checksums, vision."),
        ("Redaction:", "OffscreenCanvas masking, span-level substitution."),
        ("Local Security:", "AES-256-GCM, PBKDF2 310k, Web Crypto API."),
        ("Server Model:", "Open-weights VLM \u2014 Llama-3.2-Vision / Qwen2.5-VL via vLLM."),
        ("Server API:", "Node.js, OpenAI-compatible endpoint, guided JSON output.")]
y = 1.44
for head, body in tech:
    subhead(s, LB + 0.18, y, LW - 0.35, "\u2022 " + head, sz=10.5)
    tb(s, LB + 0.34, y + 0.2, LW - 0.5, 0.24, [{"t": body, "sz": 10, "c": INK, "line": 0.98}])
    y += 0.45

tb(s, LB + 0.18, 5.5, LW - 0.3, 0.3, [{"t": "\u2022  Methodology and process for implementation",
                                       "sz": 11.5, "b": True, "c": BLACK}])
rect(s, LB + 0.18, 5.84, LW - 0.36, 1.1, fill=GREY, line=BLACK, lw=1.0)
steps = ["1. Read screen (DOM + AX + pixels)", "2. Decide locally \u2014 act if unambiguous",
         "3. Detect & redact every personal value", "4. Verify, then transmit handles only",
         "5. Server plans one constrained action", "6. Resolve, ground, execute, confirm"]
sy = 5.94
for i in range(0, 6, 2):
    tb(s, LB + 0.3, sy, 2.9, 0.22, [{"t": steps[i], "sz": 9.5, "c": INK}])
    tb(s, LB + 3.3, sy, 2.9, 0.22, [{"t": steps[i + 1], "sz": 9.5, "c": INK}])
    sy += 0.32

# ── right box: architecture ─────────────────────────────────────────────────
tb(s, RB + 0.18, 1.0, RW - 0.3, 0.26, [{"t": "System Architecture", "sz": 11.5, "b": True, "c": HDR,
                                        "align": PP_ALIGN.CENTER}])

rect(s, RB + 0.25, 1.36, RW - 0.5, 1.62, fill=C("F7FBFF"), line=C("9DC3E6"), lw=1.0)
tb(s, RB + 0.35, 1.42, 3.0, 0.22, [{"t": "BROWSER EXTENSION (client)", "sz": 8.5, "b": True, "c": HDR}])
col = RB + 0.4
for label in ["Content Script\nreads + acts", "Service Worker\nvault + orchestration",
              "Offscreen Doc\nvision model"]:
    fbox(s, col, 1.7, 1.68, 0.54, label, CREAM, 8.5)
    col += 1.78
fbox(s, RB + 0.4, 2.34, 5.24, 0.5,
     "ONNX Runtime Web  \u00b7  WebGPU  \u2192  WASM + SIMD fallback", C("DDEBF7"), 9)

arrow(s, RB + RW / 2, 2.98, RB + RW / 2, 3.16)
rect(s, RB + 0.25, 3.16, RW - 0.5, 0.36, fill=C("FBEBE9"), line=RED, lw=1.25)
tb(s, RB + 0.25, 3.16, RW - 0.5, 0.36, [{"t": "SANITIZED CONTEXT  \u2014  ~10 KB, handles only",
                                         "sz": 9.5, "b": True, "c": RED, "align": PP_ALIGN.CENTER}],
   anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
arrow(s, RB + RW / 2, 3.52, RB + RW / 2, 3.7)

rect(s, RB + 0.25, 3.7, RW - 0.5, 1.14, fill=C("EDEBF8"), line=C("8E86F0"), lw=1.0)
tb(s, RB + 0.35, 3.76, 3.0, 0.22, [{"t": "SERVER", "sz": 8.5, "b": True, "c": C("4A42B8")}])
col = RB + 0.4
for label in ["Node.js API\nschema guard", "Open-weights VLM\nvLLM / Ollama", "PII Output Guard\nrejects literals"]:
    fbox(s, col, 4.04, 1.68, 0.6, label, C("E4E0F7"), 8.5)
    col += 1.78

rect(s, RB + 0.25, 5.0, RW - 0.5, 1.95, fill=GREY, line=BLACK, lw=1.0)
tb(s, RB + 0.38, 5.08, RW - 0.76, 0.24, [{"t": "Working prototype \u2014 verifiable today",
                                          "sz": 10, "b": True, "c": GREEN}])
proof = [("npm run eval", "PII precision 1.000 / recall 1.000"),
         ("npm run model-check", "ONNX model loads, shapes verified"),
         ("Side panel", "privacy receipt + exact bytes sent"),
         ("Live overlay", "every control labelled on any page"),
         ("Telemetry", "per-stage ms, bytes, network calls")]
py = 5.38
for cmd, what in proof:
    tb(s, RB + 0.38, py, 1.85, 0.22, [{"t": cmd, "sz": 9, "b": True, "c": HDR, "font": MONO}])
    tb(s, RB + 2.3, py, 3.3, 0.22, [{"t": what, "sz": 9, "c": INK}])
    py += 0.3

# ═══════════════════════════════════════════════════════════ 4 · FEASIBILITY
s = slide("FEASIBILITY AND VIABILITY")
rect(s, LB, 0.9, LW, 6.25, fill=WHITE, line=BLACK, lw=1.25)
rect(s, RB, 0.9, RW, 6.25, fill=WHITE, line=BLACK, lw=1.25)

diamond_header(s, LB + 0.18, 1.0, LW - 0.3, "Feasibility:")
tb(s, LB + 0.18, 1.46, LW - 0.3, 0.24, [{"t": "\u2022  Analysis of the feasibility of the idea",
                                         "sz": 10.5, "b": True, "c": BLACK}])
groups = [
    ("1. Technical Feasibility:", [
        ("Already Running End to End:", "Extension and server work today, not a mock-up."),
        ("Genuinely Small Model:", "1.2 MB \u2014 loads faster than most web fonts."),
    ]),
    ("2. Operational Feasibility:", [
        ("No New Hardware:", "Any laptop from the last five years; no GPU needed."),
        ("Works on Any Website:", "Runs on government portals and private forms alike."),
    ]),
    ("3. Legal & Ethical Feasibility:", [
        ("DPDP Act 2023 Compliance:", "Minimisation is structural \u2014 not a policy promise."),
        ("User Consent & Control:", "Irreversible actions always stop for human approval."),
    ]),
]
y = 1.76
for gname, rows in groups:
    tb(s, LB + 0.18, y, LW - 0.35, 0.24, [{"t": gname, "sz": 10.5, "b": True, "c": GREEN, "u": True}])
    y += 0.26
    for head, body in rows:
        tb(s, LB + 0.3, y, LW - 0.5, 0.22, [{"t": head, "sz": 10, "b": True, "c": INK}])
        tb(s, LB + 0.48, y + 0.2, LW - 0.7, 0.22,
           [{"t": "\u2022 " + body, "sz": 9.5, "c": INK, "line": 0.98}])
        y += 0.44
    y += 0.1

tb(s, LB + 0.18, 4.9, LW - 0.3, 0.24, [{"t": "\u2022  Sustainability of the solution",
                                        "sz": 10.5, "b": True, "c": BLACK}])
y = 5.18
for head, body in [("Open Standards Only:", "WebExtension, ONNX, WebGPU \u2014 nothing proprietary."),
                   ("No Vendor Lock-In:", "Open-weights server model, deployable offline."),
                   ("Low Running Cost:", "No GPU bill \u2014 perception runs on the client."),
                   ("Survives Model Churn:", "Swap the model file; the design is unchanged.")]:
    tb(s, LB + 0.3, y, LW - 0.5, 0.22, [{"t": head, "sz": 10, "b": True, "c": INK}])
    tb(s, LB + 0.48, y + 0.19, LW - 0.7, 0.22, [{"t": "\u2022 " + body, "sz": 9.5, "c": INK, "line": 0.98}])
    y += 0.44

diamond_header(s, RB + 0.18, 1.0, RW - 0.3, "Viability:")
tb(s, RB + 0.18, 1.46, RW - 0.3, 0.24, [{"t": "\u2022  Potential challenges and risks",
                                         "sz": 10.5, "b": True, "c": BLACK}])
tb(s, RB + 0.18, 1.7, RW - 0.3, 0.24, [{"t": "\u2022  Strategies for overcoming these challenges",
                                        "sz": 10.5, "b": True, "c": BLACK}])
risks = [("Vision too slow inside a browser tab", "Only 5\u201320 small crops per step, never the whole screen."),
         ("No WebGPU on the judging machine", "Falls back to CPU; Fast mode needs no vision at all."),
         ("Over-redacting loses marks", "Checksums prove when a number is NOT an identity number."),
         ("Page changes while the agent thinks", "Re-checks the element signature before every click."),
         ("Hostile page attempts prompt injection", "A value can only enter a field of the matching type."),
         ("A value silently fails to enter", "Every field is read back and compared after filling."),
         ("Firefox build not yet tested", "Same codebase, separate manifest \u2014 a known gap.")]
y = 2.0
for risk, fix in risks:
    tb(s, RB + 0.3, y, RW - 0.5, 0.22, [{"t": "\u25B2  " + risk, "sz": 10, "b": True, "c": RED}])
    tb(s, RB + 0.48, y + 0.2, RW - 0.7, 0.22,
       [{"t": "\u2713  " + fix, "sz": 9.5, "c": INK, "line": 0.98}])
    y += 0.46

rect(s, RB + 0.18, 5.3, RW - 0.36, 1.65, fill=GREY, line=BLACK, lw=1.0)
tb(s, RB + 0.32, 5.38, RW - 0.64, 0.24, [{"t": "Scalability & Economic Viability", "sz": 10.5,
                                          "b": True, "c": GREEN, "u": True}])
y = 5.66
for head, body in [("Cost:", "Heavy work runs on a million clients, not on rented GPUs."),
                   ("Bandwidth:", "About 10 KB per step instead of full screenshots."),
                   ("Breach risk:", "No user data stored server-side \u2014 nothing to leak."),
                   ("Portability:", "The same boundary design fits desktop and mobile agents.")]:
    tb(s, RB + 0.34, y, 1.15, 0.22, [{"t": head, "sz": 9.5, "b": True, "c": INK}])
    tb(s, RB + 1.5, y, RW - 1.75, 0.22, [{"t": body, "sz": 9.5, "c": INK}])
    y += 0.31

# ═══════════════════════════════════════════════════════════ 5 · IMPACT
s = slide("IMPACT AND BENEFITS")
rect(s, LB, 0.9, LW, 6.25, fill=WHITE, line=BLACK, lw=1.25)
rect(s, RB, 0.9, RW, 6.25, fill=WHITE, line=BLACK, lw=1.25)

diamond_header(s, LB + 0.18, 1.0, LW - 0.3, "Impacts:")
tb(s, LB + 0.18, 1.46, LW - 0.3, 0.24, [{"t": "\u2022  Potential impact on the target audience",
                                         "sz": 10.5, "b": True, "c": BLACK}])
impacts = [
    ("1. Citizens on Government Portals:", [
        "Aadhaar and PAN stay on the phone while the form is filled.",
        "Scholarship and pension applications completed without exposure.",
    ]),
    ("2. Banking, Insurance & Finance:", [
        "Card and account details never reach a third-party model.",
        "Automation becomes possible where policy previously forbade it.",
    ]),
    ("3. Healthcare & Legal Practice:", [
        "Patient and client records stay inside the office.",
        "Confidentiality obligations are met by design, not by promise.",
    ]),
    ("4. Accessibility \u2014 Blind & Low-Vision Users:", [
        "An agent reads and operates the page on the user's behalf.",
        "Real assistance without surrendering the entire screen.",
    ]),
    ("5. Every Web User:", [
        "Raises the ceiling on what people will let an agent do at all.",
        "Privacy becomes inspectable \u2014 a receipt per step, not a policy page.",
    ]),
]
y = 1.76
for gname, rows in impacts:
    tb(s, LB + 0.18, y, LW - 0.35, 0.24, [{"t": gname, "sz": 10.5, "b": True, "c": GREEN, "u": True}])
    y += 0.25
    for r in rows:
        tb(s, LB + 0.42, y, LW - 0.6, 0.22, [{"t": "\u2022 " + r, "sz": 9.8, "c": INK, "line": 0.98}])
        y += 0.24
    y += 0.12

diamond_header(s, RB + 0.18, 1.0, RW - 0.3, "Benefits:")
tb(s, RB + 0.18, 1.46, RW - 0.3, 0.24,
   [{"t": "\u2022  Benefits of the solution (social, economic, environmental, etc.)",
     "sz": 10.5, "b": True, "c": BLACK}])
bens = [
    ("1. Social Benefits:", ["People trust an agent with tasks they refuse today.",
                             "Privacy you can read on screen, not a policy nobody opens."]),
    ("2. Economic Benefits:", ["No GPU bill for perception \u2014 it runs on the client.",
                               "~10 KB per step instead of a screenshot per step."]),
    ("3. Environmental Benefits:", ["92% of each screen is work that is simply never done.",
                                    "Most steps use no network and no server compute at all."]),
    ("4. Security Benefits:", ["Passwords are removed, never masked \u2014 never read at all.",
                               "Prompt-injection blocked structurally by typed sinks."]),
    ("5. Governance Benefits:", ["Per-step receipts create an auditable privacy trail.",
                                 "DPDP data minimisation enforced by the transport format."]),
]
y = 1.76
for gname, rows in bens:
    tb(s, RB + 0.18, y, RW - 0.35, 0.24, [{"t": gname, "sz": 10.5, "b": True, "c": GREEN, "u": True}])
    y += 0.25
    for r in rows:
        tb(s, RB + 0.42, y, RW - 0.6, 0.22, [{"t": "\u2022 " + r, "sz": 9.8, "c": INK, "line": 0.98}])
        y += 0.24
    y += 0.12

rect(s, RB + 0.18, 5.35, RW - 0.36, 1.6, fill=C("E9F2EC"), line=GREEN, lw=1.0)
tb(s, RB + 0.32, 5.44, RW - 0.64, 0.24, [{"t": "National Alignment & Future Scope", "sz": 10.5,
                                          "b": True, "c": GREEN, "u": True}])
y = 5.72
for head, body in [("Digital India:", "Safe automation for citizens on public portals."),
                   ("Atmanirbhar Bharat:", "Open-weights and self-hosted \u2014 no foreign API."),
                   ("DPDP Act 2023:", "The server is structurally unable to receive the data."),
                   ("Future Scope:", "Local NER & OCR, desktop agents, a published scheme.")]:
    tb(s, RB + 0.34, y, 1.5, 0.22, [{"t": head, "sz": 9.5, "b": True, "c": INK}])
    tb(s, RB + 1.86, y, RW - 2.1, 0.22, [{"t": body, "sz": 9.5, "c": INK}])
    y += 0.3

# ═══════════════════════════════════════════════════════════ 6 · REFERENCES
s = slide("RESEARCH  AND REFERENCES")
rect(s, LB, 0.9, W - 2 * LB, 6.25, fill=WHITE, line=BLACK, lw=1.25)

tb(s, LB + 0.25, 1.0, 11.0, 0.3, [{"t": "\u2022  Details / Links of the reference and research work",
                                   "sz": 12, "b": True, "c": BLACK}])

refs = [
    ("Problem Statement:", "SIH26171 \u2014 On-device Visual Perception for Light-weight Browser Agents. "
                           "Client-side ViT or equivalent CV model, local PII redaction before any network request.",
     ["https://sih.gov.in"]),
    ("Chrome Extensions \u2014 Manifest V3:", "Platform, permissions, service workers and offscreen documents "
                                             "used to isolate the vault from the page.",
     ["https://developer.chrome.com/docs/extensions/develop/migrate",
      "https://developer.chrome.com/docs/extensions/reference/api/offscreen"]),
    ("WebGPU & ONNX Runtime Web:", "Browser GPU compute and the runtime that executes our model on-device, "
                                   "with a WASM + SIMD fallback path.",
     ["https://www.w3.org/TR/webgpu/", "https://onnxruntime.ai/docs/tutorials/web/"]),
    ("UltraFace RFB-320 \u2014 ONNX Model Zoo:", "The 1.2 MB face-detection CNN we ship and run inside the "
                                                "browser; verified input/output contract.",
     ["https://github.com/onnx/models"]),
    ("Open-Weights Server Models:", "Llama-3.2-Vision and Qwen2.5-VL served through vLLM or Ollama \u2014 "
                                    "offline deployable, as the PS requires.",
     ["https://huggingface.co", "https://docs.vllm.ai"]),
    ("Identity Checksums:", "Verhoeff for Aadhaar and Luhn (ISO/IEC 7812) for cards \u2014 how we keep PII "
                            "precision high while recall stays at 1.000.",
     ["https://uidai.gov.in"]),
    ("DPDP Act 2023 & NIST SP 800-132:", "Data minimisation obligations, and the PBKDF2 guidance behind our "
                                         "310,000-iteration key derivation.",
     ["https://www.meity.gov.in", "https://csrc.nist.gov/pubs/sp/800/132/final"]),
    ("W3C Accessible Name Computation:", "How we derive the human label of every control, which makes the "
                                         "ScreenGraph stable across site redesigns.",
     ["https://www.w3.org/TR/accname-1.2/"]),
]
y = 1.42
for i, (head, body, links) in enumerate(refs, 1):
    tb(s, LB + 0.25, y, 11.6, 0.24, [{"t": [(f"{i}. ", {"b": True, "c": BLACK}),
                                            (head, {"b": True, "c": GREEN, "u": True}),
                                            ("  " + body, {"c": INK})], "sz": 10, "line": 0.98}])
    tb(s, LB + 0.5, y + 0.24, 11.4, 0.22,
       [{"t": "   ".join(links), "sz": 9, "c": BLUE, "u": True, "line": 0.98}])
    y += 0.68

# lock-safe save
import os
BASE = r"c:\Users\srika\OneDrive\Desktop\sih-p\docs\SIH2026_Cordon_Idea"
target = f"{BASE}.pptx"
for attempt in range(1, 40):
    try:
        prs.save(target)
        break
    except PermissionError:
        target = f"{BASE}_{attempt}.pptx"
else:
    raise SystemExit("could not save — close the deck in PowerPoint and re-run")

print(f"saved {_n[0]} slides  ->  {os.path.basename(target)}")
if target != f"{BASE}.pptx":
    print("   (usual filename locked by PowerPoint — close it and re-run to overwrite)")
